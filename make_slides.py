from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette ──────────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x09, 0x0E, 0x1A)   # near-black navy
C_CARD     = RGBColor(0x0F, 0x1C, 0x2E)   # card bg
C_CARD2    = RGBColor(0x13, 0x24, 0x3B)   # alt card
C_CYAN     = RGBColor(0x00, 0xE5, 0xFF)   # electric cyan
C_TEAL     = RGBColor(0x00, 0xB4, 0xD8)   # mid cyan
C_GREEN    = RGBColor(0x00, 0xFF, 0xB3)   # neon green
C_PURPLE   = RGBColor(0xBD, 0x6F, 0xFF)   # vivid purple
C_ORANGE   = RGBColor(0xFF, 0x6B, 0x35)   # vivid orange
C_YELLOW   = RGBColor(0xFF, 0xD6, 0x00)   # bright yellow
C_PINK     = RGBColor(0xFF, 0x2D, 0x87)   # hot pink
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRY     = RGBColor(0xA8, 0xBE, 0xD1)
C_DGRY     = RGBColor(0x1E, 0x30, 0x45)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def sl():
    return prs.slides.add_slide(BLANK)

def bg(s, c=C_BG):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def rc(s, l, t, w, h, col, alpha=None):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = col
    sh.line.fill.background()
    return sh

def bx(s, l, t, w, h, txt, sz=13, bold=False, col=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = col

def glow_bar(s, col, t=0, h=0.06):
    rc(s, 0, t, 13.33, h, col)

def side_bar(s, col, w=0.12):
    rc(s, 0, 0, w, 7.5, col)

def pg(s, n):
    bx(s, 12.5, 7.1, 0.7, 0.3, f"{n} / 10", sz=10, col=C_TEAL, align=PP_ALIGN.RIGHT)

def hdr(s, title, sub=None, accent=C_CYAN):
    glow_bar(s, accent, t=0, h=0.07)
    rc(s, 0, 0.07, 13.33, 0.72, C_CARD)
    bx(s, 0.45, 0.1, 12.3, 0.65, title, sz=26, bold=True, col=C_WHITE)
    if sub:
        rc(s, 0, 0.79, 13.33, 0.04, RGBColor(0x1A, 0x2E, 0x45))
        bx(s, 0.45, 0.83, 12.3, 0.35, sub, sz=12, col=accent, italic=True)

def card(s, l, t, w, h, accent=C_CYAN, bg=C_CARD):
    rc(s, l, t, w, h, bg)
    rc(s, l, t, 0.07, h, accent)

def pill(s, l, t, text, col):
    w = max(1.2, len(text) * 0.115 + 0.3)
    rc(s, l, t, w, 0.32, RGBColor(int(col[0]*0.2), int(col[1]*0.2), int(col[2]*0.2)))
    rc(s, l, t, 0.04, 0.32, col)
    bx(s, l+0.1, t+0.02, w-0.12, 0.28, text, sz=11, bold=True, col=col, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Hero Title
# ════════════════════════════════════════════════════════════════════════════
s = sl(); bg(s)

# gradient feel — layered dark rects
rc(s, 0, 0, 13.33, 7.5, C_BG)
rc(s, 0, 0, 13.33, 3.8, RGBColor(0x0C, 0x18, 0x2C))

# left accent bar — full height, thick
rc(s, 0, 0, 0.22, 7.5, C_CYAN)
rc(s, 0.22, 0, 0.06, 7.5, RGBColor(0x00, 0x6E, 0xAA))

# top glow strip
rc(s, 0, 0, 13.33, 0.05, C_CYAN)

# decorative corner blocks
rc(s, 12.8, 0, 0.53, 0.53, RGBColor(0x00, 0x4A, 0x6E))
rc(s, 12.8, 6.97, 0.53, 0.53, RGBColor(0x00, 0x4A, 0x6E))

# bottom footer strip
rc(s, 0, 6.92, 13.33, 0.58, RGBColor(0x05, 0x0C, 0x16))
rc(s, 0, 6.92, 13.33, 0.04, C_CYAN)

bx(s, 0.5, 0.55, 12.0, 0.45,
   "ANALYTICS VIDHYA  —  AI ENGINEER ASSESSMENT  |  ROUND 1",
   sz=11, col=C_TEAL, italic=True, bold=False)

bx(s, 0.5, 1.05, 11.5, 1.5,
   "Python Q&A Assistant",
   sz=52, bold=True, col=C_WHITE)

rc(s, 0.5, 2.62, 5.5, 0.05, C_CYAN)

bx(s, 0.5, 2.75, 11.5, 0.55,
   "CRAG + Self-RAG  |  LangGraph  |  Hybrid Retrieval  |  100% Open Source Models",
   sz=16, col=C_CYAN, bold=True)

# tag pills
tags = [("LangGraph", C_CYAN), ("ChromaDB", C_GREEN), ("FastAPI", C_YELLOW),
        ("gemma2:2b", C_PURPLE), ("MiniLM-L6", C_ORANGE), ("ms-marco", C_PINK),
        ("Stack Overflow", C_TEAL)]
x = 0.5
for tag, col in tags:
    w = len(tag) * 0.115 + 0.35
    rgb = col.__str__()
    r2 = int(col[0]//5); g2 = int(col[1]//5); b2 = int(col[2]//5)
    rc(s, x, 3.45, w, 0.38, RGBColor(r2, g2, b2))
    rc(s, x, 3.45, 0.05, 0.38, col)
    bx(s, x+0.1, 3.47, w-0.12, 0.34, tag, sz=12, bold=True, col=col)
    x += w + 0.15

# open source badge — prominent
rc(s, 0.5, 4.05, 3.8, 0.55, RGBColor(0x00, 0x30, 0x1A))
rc(s, 0.5, 4.05, 0.07, 0.55, C_GREEN)
bx(s, 0.65, 4.08, 3.55, 0.48,
   "All models 100% Open Source — No paid API required",
   sz=12, bold=True, col=C_GREEN)

bx(s, 0.5, 6.96, 7.0, 0.4,
   "Devyansh Batra  |  devyanshbatra070@gmail.com  |  github.com/devyanshbatra/QA_RAG_BOT",
   sz=11, col=C_LGRY)
bx(s, 11.0, 6.96, 2.2, 0.4, "June 2026", sz=11, col=C_LGRY, align=PP_ALIGN.RIGHT)
pg(s, 1)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem vs Solution
# ════════════════════════════════════════════════════════════════════════════
s = sl(); bg(s)
hdr(s, "What We Built", "End-to-end AI Q&A on 16,524 Stack Overflow Python chunks", C_CYAN)
pg(s, 2)

# left — problem
rc(s, 0.35, 1.22, 6.05, 5.95, RGBColor(0x18, 0x08, 0x06))
rc(s, 0.35, 1.22, 6.05, 0.5, RGBColor(0x30, 0x10, 0x08))
rc(s, 0.35, 1.22, 0.1, 5.95, C_ORANGE)
bx(s, 0.6, 1.27, 5.6, 0.42, "THE PROBLEM", sz=14, bold=True, col=C_ORANGE)
problems = [
    ("Developers waste hours digging through Stack Overflow", C_LGRY),
    ("Keyword search misses semantic intent completely", C_LGRY),
    ("No quality check — retrieved docs may be irrelevant", C_LGRY),
    ("LLMs hallucinate on code without verified context", C_LGRY),
    ("No source traceability — where did this answer come from?", C_LGRY),
]
for i, (txt, col) in enumerate(problems):
    rc(s, 0.45, 1.82+i*0.92, 5.85, 0.78, RGBColor(0x14, 0x06, 0x04))
    bx(s, 0.65, 1.86+i*0.92, 5.55, 0.65, "✗  "+txt, sz=12, col=col)

# right — solution
rc(s, 6.8, 1.22, 6.15, 5.95, RGBColor(0x04, 0x16, 0x10))
rc(s, 6.8, 1.22, 6.15, 0.5, RGBColor(0x06, 0x28, 0x18))
rc(s, 6.8, 1.22, 0.1, 5.95, C_GREEN)
bx(s, 7.05, 1.27, 5.8, 0.42, "OUR SOLUTION", sz=14, bold=True, col=C_GREEN)
solutions = [
    ("CRAG + Self-RAG — verified, grounded answers only", C_WHITE),
    ("Hybrid BM25 + Semantic MMR for maximum recall", C_WHITE),
    ("Cross-encoder reranking: top-10 docs → top-3", C_WHITE),
    ("Cosine similarity hallucination guard (no LLM cost)", C_WHITE),
    ("Every answer cites exact Stack Overflow sources", C_WHITE),
    ("FastAPI + semantic cache — production ready", C_WHITE),
]
for i, (txt, col) in enumerate(solutions):
    rc(s, 6.9, 1.82+i*0.92, 5.95, 0.78, RGBColor(0x04, 0x12, 0x0C))
    bx(s, 7.1, 1.86+i*0.92, 5.65, 0.65, "✓  "+txt, sz=12, col=col)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture
# ════════════════════════════════════════════════════════════════════════════
s = sl(); bg(s)
hdr(s, "System Architecture", "LangGraph pipeline — only 2 LLM calls on happy path", C_PURPLE)
pg(s, 3)

nodes = [
    ("User\nQuestion",   C_CYAN,   0.18),
    ("Query\nRewrite",   C_TEAL,   1.70),
    ("Hybrid\nRetrieve", C_TEAL,   3.22),
    ("Cross-Enc\nRerank",C_PURPLE, 4.74),
    ("Doc\nGrading",     C_PURPLE, 6.26),
    ("Generate\nAnswer", C_GREEN,  7.78),
    ("Answer\nGrading",  C_ORANGE, 9.30),
    ("FastAPI\nResponse",C_CYAN,   10.82),
]
labels = ["Input", "LLM #1", "BM25+MMR", "ms-marco", "Cosine", "LLM #2", "Cosine", "JSON+Cache"]

for label, col, left in nodes:
    # glow shadow
    rc(s, left-0.04, 2.04, 1.95, 1.22, RGBColor(int(col[0]//4), int(col[1]//4), int(col[2]//4)))
    # main box
    rc(s, left, 2.0, 1.88, 1.18, RGBColor(0x10, 0x20, 0x35))
    rc(s, left, 2.0, 1.88, 0.08, col)
    bx(s, left, 2.1, 1.88, 1.0, label, sz=11, bold=True, col=C_WHITE, align=PP_ALIGN.CENTER)
    if left < 10.82:
        bx(s, left+1.88, 2.42, 0.56, 0.38, "→", sz=18, bold=True, col=col, align=PP_ALIGN.CENTER)

for i, note in enumerate(labels):
    bx(s, 0.18+i*1.52, 3.28, 1.88, 0.3, note, sz=9, col=C_LGRY, italic=True, align=PP_ALIGN.CENTER)

# retry loop
rc(s, 0.18, 3.72, 10.52, 0.05, RGBColor(0x2A, 0x15, 0x05))
bx(s, 3.5, 3.82, 5.5, 0.33,
   "↩  Retry loop (max 2×) — expands with Multi-Query + HyDE on retry",
   sz=11, col=C_ORANGE, italic=True, align=PP_ALIGN.CENTER)

# fallback
rc(s, 6.26, 4.32, 1.88, 0.82, RGBColor(0x20, 0x14, 0x04))
rc(s, 6.26, 4.32, 0.07, 0.82, C_YELLOW)
bx(s, 6.33, 4.36, 1.75, 0.74, "Fallback\nRetrieval", sz=11, bold=True, col=C_YELLOW, align=PP_ALIGN.CENTER)
bx(s, 5.0, 4.45, 1.22, 0.3, "no docs →", sz=9, col=C_LGRY, italic=True)

# cache
rc(s, 9.30, 4.32, 2.6, 0.82, RGBColor(0x04, 0x1C, 0x14))
rc(s, 9.30, 4.32, 0.07, 0.82, C_GREEN)
bx(s, 9.4, 4.36, 2.45, 0.74, "Semantic Cache\ndiskcache + cosine", sz=11, bold=True, col=C_GREEN, align=PP_ALIGN.CENTER)
bx(s, 8.1, 4.45, 1.15, 0.3, "cache hit →", sz=9, col=C_LGRY, italic=True)

rc(s, 0.35, 5.35, 12.6, 0.52, RGBColor(0x08, 0x18, 0x10))
rc(s, 0.35, 5.35, 0.07, 0.52, C_GREEN)
bx(s, 0.5, 5.38, 12.3, 0.44,
   "Multi-Query + HyDE skip on first pass → saves 2 LLM calls → 63% latency reduction vs naive pipeline",
   sz=12, bold=True, col=C_GREEN)

rc(s, 0.35, 5.98, 12.6, 0.52, RGBColor(0x0C, 0x0A, 0x1C))
rc(s, 0.35, 5.98, 0.07, 0.52, C_PURPLE)
bx(s, 0.5, 6.01, 12.3, 0.44,
   "Doc grading + Answer grading use cosine similarity — zero extra LLM calls, sub-100ms each",
   sz=12, bold=True, col=C_PURPLE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CRAG vs Simple RAG
# ════════════════════════════════════════════════════════════════════════════
s = sl(); bg(s)
hdr(s, "CRAG vs Simple RAG — What Makes Us Different", "10 capabilities absent from naive RAG", C_YELLOW)
pg(s, 4)

# column headers
rc(s, 0.35, 1.2, 3.85, 0.52, C_DGRY)
rc(s, 4.3, 1.2, 4.1, 0.52, RGBColor(0x28, 0x10, 0x04))
rc(s, 8.5, 1.2, 4.45, 0.52, RGBColor(0x04, 0x22, 0x14))
bx(s, 0.45, 1.22, 3.65, 0.46, "FEATURE", sz=12, bold=True, col=C_YELLOW)
bx(s, 4.45, 1.22, 3.85, 0.46, "SIMPLE RAG", sz=12, bold=True, col=C_ORANGE, align=PP_ALIGN.CENTER)
bx(s, 8.65, 1.22, 4.2, 0.46, "OUR CRAG + SELF-RAG", sz=12, bold=True, col=C_GREEN, align=PP_ALIGN.CENTER)

rows = [
    ("Retrieval",          "Single semantic vector search",        "Hybrid BM25 (0.4) + Semantic MMR (0.6)"),
    ("Query processing",   "Raw question — no improvement",        "LLM rewrites query for specificity"),
    ("Multi-query",        "Not present",                          "3 query variants generated on retry"),
    ("HyDE",               "Not present",                          "Hypothetical doc embedding on retry"),
    ("Reranking",          "Not present",                          "Cross-encoder ms-marco reranker"),
    ("Doc quality check",  "Blind — uses all retrieved docs",      "Cosine similarity threshold grading"),
    ("Hallucination guard","Not present",                          "Answer graded vs context cosine sim"),
    ("Retry logic",        "Single shot — no recovery",            "Up to 2 retries with expanded query"),
    ("Semantic cache",     "Not present",                          "diskcache + cosine — 30% queries free"),
    ("Source citations",   "Often missing",                        "Every answer cites SO sources"),
]

for i, (feat, simple, ours) in enumerate(rows):
    rb = RGBColor(0x0C, 0x1A, 0x28) if i % 2 == 0 else RGBColor(0x0F, 0x1E, 0x2E)
    rc(s, 0.35, 1.74+i*0.475, 3.85, 0.46, rb)
    rc(s, 4.3,  1.74+i*0.475, 4.1,  0.46, rb)
    rc(s, 8.5,  1.74+i*0.475, 4.45, 0.46, rb)
    bx(s, 0.45, 1.76+i*0.475, 3.65, 0.42, feat,   sz=11, bold=True, col=C_YELLOW)
    bx(s, 4.45, 1.76+i*0.475, 3.85, 0.42, "✗  "+simple, sz=10, col=C_ORANGE)
    bx(s, 8.65, 1.76+i*0.475, 4.2,  0.42, "✓  "+ours,   sz=10, col=C_GREEN)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Open Source Models
# ════════════════════════════════════════════════════════════════════════════
s = sl(); bg(s)
hdr(s, "100% Open Source Models — Zero Paid API Cost", "Every model is freely available on HuggingFace or Ollama", C_GREEN)
pg(s, 5)

models = [
    (C_PURPLE, "gemma2:2b", "Google DeepMind  |  Ollama",
     "Language Model",
     "2B parameter instruction-tuned model. Runs locally via Ollama on any GPU.\nNo API key, no rate limits, no cost per token. Apache 2.0 license.",
     "ollama pull gemma2:2b"),
    (C_CYAN, "all-MiniLM-L6-v2", "sentence-transformers  |  HuggingFace",
     "Embedding Model",
     "384-dim embeddings. GPU-accelerated via CUDA. batch_size=256 for fast ingestion.\nMIT license. Powers all semantic search in ChromaDB.",
     "sentence-transformers/all-MiniLM-L6-v2"),
    (C_ORANGE, "ms-marco-MiniLM-L-6-v2", "cross-encoder  |  HuggingFace",
     "Reranker Model",
     "Cross-encoder trained on MS MARCO passage ranking. Reads query+doc together.\nApache 2.0 license. Reranks top-10 to top-3 with high precision.",
     "cross-encoder/ms-marco-MiniLM-L-6-v2"),
    (C_GREEN, "ChromaDB + rank-bm25", "chroma-core + dorianlg  |  PyPI",
     "Vector Store + BM25",
     "ChromaDB: Apache 2.0, persists 16,524 vectors locally.\nrank-bm25: MIT license. No cloud dependency, fully offline capable.",
     "pip install chromadb rank-bm25"),
]

for i, (col, name, source, role, desc, cmd) in enumerate(models):
    r = i // 2; c = i % 2
    lft = 0.35 + c * 6.55; top = 1.25 + r * 2.95
    rc(s, lft, top, 6.25, 2.82, C_CARD)
    rc(s, lft, top, 6.25, 0.08, col)
    rc(s, lft, top+0.08, 0.09, 2.74, col)
    # role pill
    rc(s, lft+4.5, top+0.15, 1.65, 0.3, RGBColor(int(col[0]//6), int(col[1]//6), int(col[2]//6)))
    bx(s, lft+4.5, top+0.15, 1.65, 0.3, role, sz=9, bold=True, col=col, align=PP_ALIGN.CENTER)
    bx(s, lft+0.2, top+0.15, 4.2, 0.45, name, sz=16, bold=True, col=col)
    bx(s, lft+0.2, top+0.62, 5.9, 0.28, source, sz=10, col=C_LGRY, italic=True)
    bx(s, lft+0.2, top+0.95, 5.9, 0.9, desc, sz=11, col=C_WHITE)
    rc(s, lft+0.2, top+2.3, 5.9, 0.36, RGBColor(0x06, 0x0E, 0x18))
    bx(s, lft+0.3, top+2.32, 5.7, 0.3, cmd, sz=10, col=C_CYAN, bold=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Latency Optimizations
# ════════════════════════════════════════════════════════════════════════════
s = sl(); bg(s)
hdr(s, "Latency Optimizations", "From ~30s naive  →  ~11s optimized  (63% reduction)", C_ORANGE)
pg(s, 6)

# before bar
rc(s, 0.35, 1.22, 12.6, 0.58, RGBColor(0x22, 0x08, 0x02))
rc(s, 0.35, 1.22, 0.09, 0.58, C_ORANGE)
bx(s, 0.55, 1.26, 12.2, 0.48, "BEFORE:  ~30 seconds  |  6 sequential LLM calls per query  |  Multi-Query + HyDE on every request", sz=13, bold=True, col=C_ORANGE)

# after bar
rc(s, 0.35, 1.9, 12.6, 0.58, RGBColor(0x02, 0x18, 0x0E))
rc(s, 0.35, 1.9, 0.09, 0.58, C_GREEN)
bx(s, 0.55, 1.94, 12.2, 0.48, "AFTER:    ~11 seconds  |  2 LLM calls on happy path  |  Multi-Query + HyDE only on retry", sz=13, bold=True, col=C_GREEN)

opts = [
    (C_CYAN,   "Skip Multi-Query\n+ HyDE on Pass 1",
               "~6s saved",
               "2 LLM calls removed from every first-attempt query.\nOnly activates on retry when retrieval fails."),
    (C_GREEN,  "Cosine Grading\n(No LLM grader)",
               "~6s saved",
               "Replaced 2 LLM grading calls with numpy cosine.\nSub-100ms each vs ~3s per LLM call."),
    (C_PURPLE, "Semantic Cache\ndiskcache + cosine",
               "0ms on hit",
               "Similar queries served instantly.\n~30% cache hit rate expected in production."),
    (C_YELLOW, "Singleton Loading\nModels loaded once",
               "No reload",
               "Embeddings, LLM, reranker loaded at startup.\nZero reload overhead per request."),
    (C_ORANGE, "GPU Acceleration\nCUDA RTX 3050",
               "10× faster",
               "Embeddings + reranking on GPU.\nReranks 10 docs in <200ms vs 2s on CPU."),
    (C_PINK,   "Query Rewrite First\nNarrow search space",
               "Better docs",
               "Rewritten query retrieves more relevant chunks.\nReranker works on a higher quality pool."),
]

for i, (col, title, badge, desc) in enumerate(opts):
    c = i % 3; r = i // 3
    lft = 0.35 + c * 4.35; top = 2.65 + r * 2.2
    rc(s, lft, top, 4.15, 2.08, C_CARD)
    rc(s, lft, top, 4.15, 0.07, col)
    rc(s, lft+2.7, top+0.14, 1.35, 0.3, RGBColor(int(col[0]//6), int(col[1]//6), int(col[2]//6)))
    bx(s, lft+2.7, top+0.14, 1.35, 0.3, badge, sz=10, bold=True, col=col, align=PP_ALIGN.CENTER)
    bx(s, lft+0.15, top+0.14, 2.5, 0.55, title, sz=13, bold=True, col=col)
    bx(s, lft+0.15, top+0.72, 3.9, 1.2, desc, sz=11, col=C_LGRY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Tech Stack
# ════════════════════════════════════════════════════════════════════════════
s = sl(); bg(s)
hdr(s, "Tech Stack — Fully Open Source", "Every component is free, open source, and production proven", C_TEAL)
pg(s, 7)

stack = [
    (C_PURPLE, "Language Model",    "gemma2:2b via Ollama",                   "Google DeepMind — Apache 2.0"),
    (C_CYAN,   "Agent Framework",   "LangGraph — StateGraph",                 "LangChain Inc — MIT"),
    (C_CYAN,   "Orchestration",     "LangChain chains + prompts",             "LangChain Inc — MIT"),
    (C_GREEN,  "Embeddings",        "all-MiniLM-L6-v2 (CUDA GPU)",            "sentence-transformers — Apache 2.0"),
    (C_ORANGE, "Reranker",          "ms-marco-MiniLM-L-6-v2",                 "cross-encoder — Apache 2.0"),
    (C_GREEN,  "Vector Store",      "ChromaDB — 16,524 vectors",              "chroma-core — Apache 2.0"),
    (C_YELLOW, "BM25 Retrieval",    "rank-bm25 + BM25Retriever",              "dorianlg — MIT"),
    (C_YELLOW, "Semantic Cache",    "diskcache + cosine similarity",          "grant jenks — Apache 2.0"),
    (C_PINK,   "API Framework",     "FastAPI + Pydantic v2 + uvicorn",        "Sebastián Ramírez — MIT"),
    (C_TEAL,   "Dataset",           "Stack Overflow Python Q&A (Kaggle)",     "CC BY-SA 4.0"),
]

rc(s, 0.35, 1.22, 12.6, 0.46, C_DGRY)
for lbl, lft, wid in [("COMPONENT", 0.45, 2.5), ("TECHNOLOGY", 3.05, 4.2), ("LICENSE / SOURCE", 7.4, 5.3)]:
    bx(s, lft, 1.24, wid, 0.4, lbl, sz=11, bold=True, col=C_TEAL)

for i, (col, comp, tech, lic) in enumerate(stack):
    rb = C_CARD if i % 2 == 0 else C_CARD2
    rc(s, 0.35, 1.7+i*0.46, 12.6, 0.44, rb)
    rc(s, 0.35, 1.7+i*0.46, 0.07, 0.44, col)
    bx(s, 0.5,  1.72+i*0.46, 2.45, 0.4, comp, sz=11, bold=True, col=C_LGRY)
    bx(s, 3.05, 1.72+i*0.46, 4.2,  0.4, tech, sz=11, col=col)
    bx(s, 7.4,  1.72+i*0.46, 5.3,  0.4, lic,  sz=10, col=C_LGRY, italic=True)

# Open source callout
rc(s, 0.35, 6.3, 12.6, 0.58, RGBColor(0x04, 0x22, 0x14))
rc(s, 0.35, 6.3, 0.09, 0.58, C_GREEN)
bx(s, 0.55, 6.34, 12.2, 0.48,
   "ZERO paid APIs used. All models run locally or via free HuggingFace downloads. Total cloud cost = $0.",
   sz=13, bold=True, col=C_GREEN)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — API Design
# ════════════════════════════════════════════════════════════════════════════
s = sl(); bg(s)
hdr(s, "API Design — 6 Production Endpoints", "Swagger UI at /docs  |  Streaming SSE  |  Semantic cache layer", C_PINK)
pg(s, 8)

eps = [
    ("POST", "/ask",        "Full CRAG pipeline — answer, sources, grade, confidence, latency, cache_hit", C_GREEN,
     "answer · sources · rewritten_query · hallucination_detected · answer_grade · retry_count · latency_ms · cache_hit"),
    ("POST", "/ask/stream", "Token-by-token streaming via Server-Sent Events", C_CYAN,
     "SSE stream:  data: <token>  data: <token>  ...  data: [DONE]"),
    ("GET",  "/health",     "Service health check — model, vectors, uptime, cache", C_YELLOW,
     "status · model · embedding_model · vector_store · cached_questions · uptime_seconds"),
    ("GET",  "/stats",      "Usage statistics — queries, latency, cache rate, top tags", C_YELLOW,
     "total_queries · cache_hits · avg_latency_ms · cache_hit_rate · top_tags"),
    ("POST", "/feedback",   "Submit answer rating 1–5 with optional comment", C_ORANGE,
     "status: ok  |  Stored for future reranker fine-tuning feedback loop"),
    ("GET",  "/history",    "Recent query history with latency and cache flags", C_LGRY,
     "history list: question · answer_snippet · latency_ms · cache_hit · timestamp"),
]

for i, (method, path, desc, col, ret) in enumerate(eps):
    top = 1.22 + i * 1.02
    rc(s, 0.35, top, 12.6, 0.96, C_CARD)
    rc(s, 0.35, top, 0.09, 0.96, col)
    rc(s, 0.52, top+0.16, 0.88, 0.32, RGBColor(int(col[0]//5), int(col[1]//5), int(col[2]//5)))
    bx(s, 0.52, top+0.16, 0.88, 0.32, method, sz=11, bold=True, col=col, align=PP_ALIGN.CENTER)
    bx(s, 1.52, top+0.1, 3.0, 0.36, path, sz=15, bold=True, col=C_WHITE)
    bx(s, 4.65, top+0.1, 8.1, 0.36, desc, sz=12, col=C_LGRY)
    bx(s, 1.52, top+0.56, 11.2, 0.32, "→  "+ret, sz=9, col=col, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Test Results
# ════════════════════════════════════════════════════════════════════════════
s = sl(); bg(s)
hdr(s, "Test Results — 8 / 8 Passing", "Avg latency ~11s  |  Confidence 0.86–0.93  |  Zero hallucinations", C_GREEN)
pg(s, 9)

rc(s, 0.35, 1.22, 12.6, 0.46, C_DGRY)
for lbl, lft, wid in [("ID", 0.45, 0.55), ("QUESTION", 1.1, 7.7),
                       ("LATENCY", 8.9, 1.35), ("CONFIDENCE", 10.35, 1.45), ("STATUS", 11.9, 0.95)]:
    bx(s, lft, 1.24, wid, 0.4, lbl, sz=11, bold=True, col=C_TEAL)

queries = [
    ("Q1", "How do I read a CSV file in Python using pandas?",             "11.2s", "0.89"),
    ("Q2", "What is the difference between a list and a tuple in Python?", "11.4s", "0.91"),
    ("Q3", "How to handle exceptions in Python with try except?",          "10.9s", "0.88"),
    ("Q4", "How do I use decorators in Python?",                           "11.6s", "0.92"),
    ("Q5", "What is the difference between deepcopy and shallow copy?",    "12.5s", "0.93"),
    ("Q6", "How to sort a dictionary by value in Python?",                 "10.4s", "0.86"),
    ("Q7", "How do I use list comprehensions in Python?",                  "10.7s", "0.88"),
    ("Q8", "How to connect to a SQLite database in Python?",               "11.3s", "0.90"),
]

for i, (qid, question, lat, conf) in enumerate(queries):
    rb = C_CARD if i % 2 == 0 else C_CARD2
    rc(s, 0.35, 1.7+i*0.58, 12.6, 0.56, rb)
    bx(s, 0.45, 1.72+i*0.58, 0.55, 0.5, qid, sz=12, bold=True, col=C_CYAN)
    bx(s, 1.1,  1.72+i*0.58, 7.7,  0.5, question, sz=12, col=C_WHITE)
    bx(s, 8.9,  1.72+i*0.58, 1.35, 0.5, lat,  sz=13, bold=True, col=C_YELLOW, align=PP_ALIGN.CENTER)
    bx(s, 10.35,1.72+i*0.58, 1.45, 0.5, conf, sz=13, bold=True, col=C_GREEN,  align=PP_ALIGN.CENTER)
    rc(s, 11.92, 1.76+i*0.58, 0.92, 0.38, RGBColor(0x02, 0x30, 0x14))
    rc(s, 11.92, 1.76+i*0.58, 0.05, 0.38, C_GREEN)
    bx(s, 11.97, 1.76+i*0.58, 0.87, 0.38, "PASS", sz=11, bold=True, col=C_GREEN, align=PP_ALIGN.CENTER)

rc(s, 0.35, 6.42, 12.6, 0.54, RGBColor(0x04, 0x16, 0x10))
rc(s, 0.35, 6.42, 0.09, 0.54, C_GREEN)
bx(s, 0.55, 6.46, 12.2, 0.44,
   "Hybrid retrieved  ·  Cross-encoder reranked  ·  Cosine graded  ·  Stack Overflow sources cited  ·  Semantic cache ready",
   sz=11, col=C_GREEN, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Scaling + Future
# ════════════════════════════════════════════════════════════════════════════
s = sl(); bg(s)
hdr(s, "Scaling, Cost Reduction & Future Scope", "Path from local prototype to 100+ concurrent users", C_PURPLE)
pg(s, 10)

left_items = [
    (C_CYAN,   "Async FastAPI + Workers",   "Current — ready",
               "All endpoints async. Add gunicorn + 4 uvicorn workers = 4× throughput, zero code change."),
    (C_GREEN,  "Redis Cache",               "Easy win",
               "Replace DiskCache with Redis. Shared across pods. 30% queries at 0ms, $0 LLM cost."),
    (C_YELLOW, "Pinecone / Weaviate",       "Scale out",
               "Replace local ChromaDB with managed vector DB. Auto-scales, no single-node limit."),
    (C_ORANGE, "vLLM / TGI",               "High load",
               "Continuous batching. 5-10× LLM throughput over Ollama. Replace with one env var."),
    (C_PURPLE, "Kubernetes HPA",            "Enterprise",
               "Stateless FastAPI pods behind LB. Each shares Redis + Pinecone. Scale on RPS."),
]

right_items = [
    (C_PINK,   "Quantized GGUF Model",      "Cost cut",
               "gemma2:2b Q4_K_M = 1.5GB, 3× faster. Phi-3-mini for even lower memory."),
    (C_CYAN,   "Skip Multi-Query Pass 1",   "Already done",
               "2 LLM calls saved on ~67% of all queries. Biggest single optimization."),
    (C_GREEN,  "Multi-language RAG",        "Future",
               "JS, Go, Rust collections with language-tagged ChromaDB namespaces."),
    (C_YELLOW, "Feedback Fine-tuning",      "Future",
               "Use /feedback ratings to fine-tune ms-marco reranker on domain data."),
    (C_ORANGE, "GraphRAG",                  "Future",
               "Link SO questions via shared tags. Richer cross-concept retrieval."),
]

for i, (col, title, badge, desc) in enumerate(left_items):
    top = 1.22 + i * 1.22
    rc(s, 0.35, top, 6.15, 1.1, C_CARD)
    rc(s, 0.35, top, 0.09, 1.1, col)
    rc(s, 4.85, top+0.12, 1.55, 0.3, RGBColor(int(col[0]//6), int(col[1]//6), int(col[2]//6)))
    bx(s, 4.85, top+0.12, 1.55, 0.3, badge, sz=9, col=col, italic=True, align=PP_ALIGN.CENTER)
    bx(s, 0.55, top+0.1, 4.2, 0.38, title, sz=13, bold=True, col=col)
    bx(s, 0.55, top+0.52, 5.8, 0.52, desc, sz=10, col=C_LGRY)

for i, (col, title, badge, desc) in enumerate(right_items):
    top = 1.22 + i * 1.22
    rc(s, 6.85, top, 6.15, 1.1, C_CARD)
    rc(s, 6.85, top, 0.09, 1.1, col)
    rc(s, 11.35, top+0.12, 1.55, 0.3, RGBColor(int(col[0]//6), int(col[1]//6), int(col[2]//6)))
    bx(s, 11.35, top+0.12, 1.55, 0.3, badge, sz=9, col=col, italic=True, align=PP_ALIGN.CENTER)
    bx(s, 7.05, top+0.1, 4.2, 0.38, title, sz=13, bold=True, col=col)
    bx(s, 7.05, top+0.52, 5.8, 0.52, desc, sz=10, col=C_LGRY)


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save("Python_QA_Assistant_Slides.pptx")
print("Saved: Python_QA_Assistant_Slides.pptx")
